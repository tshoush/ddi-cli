import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define Luxury Colors
    BLUE_DARK = RGBColor(10, 25, 47)    # #0a192f
    GOLD_ACCENT = RGBColor(212, 175, 55) # #D4AF37
    WHITE_TEXT = RGBColor(230, 241, 255) # #e6f1ff

    def style_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BLUE_DARK

    def style_title(shape):
        shape.text_frame.paragraphs[0].font.color.rgb = GOLD_ACCENT
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.name = 'Arial'

    def style_body(shape):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE_TEXT
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(20)

    # Title Slide
    slide_layout = prs.slide_layouts[0] # 0 is title
    slide = prs.slides.add_slide(slide_layout)
    style_slide_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "DDI CLI Tool"
    style_title(title)
    title.text_frame.paragraphs[0].font.size = Pt(54)
    
    subtitle.text = "Cloud to Infoblox Synchronization"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE_TEXT
    subtitle.text_frame.paragraphs[0].font.name = 'Arial'

    # Helper to add a bullet slide
    def add_bullet_slide(heading, bullet_points):
        slide_layout = prs.slide_layouts[1] # 1 is title + content
        slide = prs.slides.add_slide(slide_layout)
        style_slide_background(slide)
        
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = heading
        style_title(title_shape)
        
        tf = body_shape.text_frame
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.color.rgb = WHITE_TEXT
            p.font.name = 'Arial'
            p.font.size = Pt(24)
            p.space_after = Pt(14)

    # Overview
    add_bullet_slide("Overview", [
        "DDI CLI is a command-line tool to sync network data from cloud providers (AWS) to Infoblox.",
        "Key Features:",
        "- Automated Sync: Keep Infoblox up-to-date.",
        "- Extensible: Modular design for new providers.",
        "- Interactive: User-friendly menu system.",
        "- Audit & Search: Tools to find and verify resources."
    ])

    # Architecture
    add_bullet_slide("Architecture", [
        "Layered architecture for flexibility:",
        "CLI Layer (ddi/cli.py): Handles user interaction (Click).",
        "Configuration (ddi/config.py): Manages settings.",
        "Core Logic (ddi/infoblox.py): Infoblox WAPI interactions.",
        "Providers (ddi/providers/): AWS implementation and base classes."
    ])

    # Installation & Setup
    add_bullet_slide("Installation & Setup", [
        "1. Clone the Repository",
        "2. Run ./setup.sh",
        "   - Creates virtual environment",
        "   - Installs dependencies",
        "3. Activate: source venv/bin/activate"
    ])

    # Interactive Workflow
    add_bullet_slide("Interactive Workflow", [
        "Configuration Dashboard:",
        "- Visual 'box' UI for settings.",
        "Network View Selection:",
        "- Dynamically fetches views from Infoblox.",
        "Main Menu:",
        "- Easy navigation to Audit, Search, and AWS commands."
    ])

    # Key Commands
    add_bullet_slide("Key Commands", [
        "audit: Scan all providers and report resources.",
        "search <term>: Find a resource across clouds.",
        "aws sync: Synchronize AWS VPC data.",
        "aws attributes analyze: Compare AWS tags with Infoblox EAs."
    ])

    # Future Roadmap
    add_bullet_slide("Future Roadmap", [
        "Azure Support: Microsoft Azure integration.",
        "GCP Support: Google Cloud Platform support.",
        "Automated Scheduling: Built-in cron/scheduler.",
        "Enhanced Reporting: HTML/PDF reports."
    ])

    # Thank You
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    style_slide_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You!"
    style_title(title)
    
    subtitle.text = "Repository: https://github.com/tshoush/ddi-cli"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE_TEXT

    prs.save('DDI_CLI_Presentation.pptx')
    print("Presentation saved to DDI_CLI_Presentation.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx is not installed.")
        print("Please run: pip install python-pptx")
